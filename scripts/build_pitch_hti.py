"""Generate HexaBin_HTI_Pitch.pptx — Armenian grant pitch deck for ԲՏԱՆ-ԴՄ-2026/01.

Run: python scripts/build_pitch_hti.py
Output: HexaBin_HTI_Pitch.pptx at repo root.

Content aligned with Առաջարկ.docx (authoritative proposal). Brand palette from CLAUDE.md.
All layout uses native python-pptx shapes — no external images.
"""
from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Inches, Pt

OUTPUT_PATH = Path(__file__).resolve().parent.parent / "HexaBin_HTI_Pitch.pptx"
TOTAL_SLIDES = 16


# ──────────────────────────── Theme ────────────────────────────


@dataclass(frozen=True)
class Theme:
    # Brand primaries
    FOREST: RGBColor = RGBColor(0x2D, 0x5A, 0x42)
    DEEP_BLUE: RGBColor = RGBColor(0x1A, 0x4D, 0x6B)
    TAUPE: RGBColor = RGBColor(0xBD, 0xB7, 0x6B)
    STONE: RGBColor = RGBColor(0x8C, 0x8C, 0x8C)

    # Module category chips
    PAPER: RGBColor = RGBColor(0xD2, 0xB4, 0x8C)
    ALUMINUM: RGBColor = RGBColor(0xA9, 0xA9, 0xA9)
    ORGANIC: RGBColor = RGBColor(0x1E, 0x4D, 0x2B)
    GLASS: RGBColor = RGBColor(0x40, 0xE0, 0xD0)
    PLASTIC: RGBColor = RGBColor(0x87, 0xCE, 0xEB)
    OTHER: RGBColor = RGBColor(0x93, 0x70, 0xDB)

    # Semantic
    SUCCESS: RGBColor = RGBColor(0x4C, 0xAF, 0x50)
    WARNING: RGBColor = RGBColor(0xFF, 0x98, 0x00)
    ERROR: RGBColor = RGBColor(0xC6, 0x28, 0x28)
    INFO: RGBColor = RGBColor(0x21, 0x96, 0xF3)

    # Neutrals
    WHITE: RGBColor = RGBColor(0xFF, 0xFF, 0xFF)
    OFFWHITE: RGBColor = RGBColor(0xF5, 0xF5, 0xF7)
    LIGHT_GRAY: RGBColor = RGBColor(0xE0, 0xE0, 0xE0)
    CHARCOAL: RGBColor = RGBColor(0x33, 0x33, 0x33)

    # Typography — Calibri ships with Office and has full Armenian coverage.
    FONT: str = "Calibri"


T = Theme()


# ──────────────────────────── Helpers ────────────────────────────


def set_fill(shape, rgb: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def set_no_fill(shape) -> None:
    shape.fill.background()


def set_no_line(shape) -> None:
    shape.line.fill.background()


def set_line(shape, rgb: RGBColor, pt: float = 1.0) -> None:
    shape.line.color.rgb = rgb
    shape.line.width = Pt(pt)


def style_text(
    tf,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = T.CHARCOAL,
    align: int = PP_ALIGN.LEFT,
    anchor: int = MSO_ANCHOR.TOP,
    font: str = T.FONT,
    bullet: bool = False,
) -> None:
    """Set a text frame to a single styled line (clears existing paragraphs)."""
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = ("• " + text) if bullet else text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_paragraph(
    tf,
    text: str,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = T.CHARCOAL,
    align: int = PP_ALIGN.LEFT,
    font: str = T.FONT,
    space_before: int = 0,
    bullet: bool = False,
) -> None:
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = ("• " + text) if bullet else text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def slide_background(slide: Slide, rgb: RGBColor) -> None:
    """Paint a full-bleed background rectangle (python-pptx can't set slide bg directly)."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    set_fill(bg, rgb)
    set_no_line(bg)
    # Send to back by reordering XML
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_title_bar(slide: Slide, title: str, number: int) -> None:
    """Deep-blue band across the top with title (left) + HexaBin · N/14 (right)."""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.9))
    set_fill(bar, T.DEEP_BLUE)
    set_no_line(bar)

    # Left accent slab (forest green) for visual rhythm
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), Inches(0.9)
    )
    set_fill(accent, T.FOREST)
    set_no_line(accent)

    tb = slide.shapes.add_textbox(Inches(0.55), Inches(0.15), Inches(9.5), Inches(0.6))
    style_text(
        tb.text_frame,
        title,
        size=26,
        bold=True,
        color=T.WHITE,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Right side: HEXABIN wordmark + page number
    right = slide.shapes.add_textbox(
        Inches(10.2), Inches(0.15), Inches(3.0), Inches(0.6)
    )
    style_text(
        right.text_frame,
        f"HEXABIN   ·   {number:02d} / {TOTAL_SLIDES:02d}",
        size=11,
        bold=True,
        color=T.WHITE,
        align=PP_ALIGN.RIGHT,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_footer(slide: Slide) -> None:
    line = slide.shapes.add_connector(1, Inches(0.5), Inches(7.1), Inches(12.83), Inches(7.1))
    line.line.color.rgb = T.LIGHT_GRAY
    line.line.width = Pt(0.75)

    left = slide.shapes.add_textbox(Inches(0.5), Inches(7.15), Inches(8.0), Inches(0.3))
    style_text(
        left.text_frame,
        "«ԷՄ ՔՅՈՒԲ» ՍՊԸ   ·   HexaBin   ·   ԲՏԱՆ-ԴՄ-2026/01",
        size=9,
        color=T.STONE,
    )
    right = slide.shapes.add_textbox(Inches(9.0), Inches(7.15), Inches(3.83), Inches(0.3))
    style_text(
        right.text_frame,
        "armovs.yan@gmail.com   ·   +374 95 851 561",
        size=9,
        color=T.STONE,
        align=PP_ALIGN.RIGHT,
    )


def add_card(
    slide: Slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    title: str,
    body: str | None = None,
    accent: RGBColor = T.FOREST,
    title_size: int = 16,
    body_size: int = 12,
) -> None:
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    card.adjustments[0] = 0.08
    set_fill(card, T.WHITE)
    set_line(card, T.LIGHT_GRAY, 1)

    # Left accent stripe
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y + 0.1), Inches(0.12), Inches(h - 0.2)
    )
    set_fill(stripe, accent)
    set_no_line(stripe)

    # Title + body text inside a textbox over the card
    tb = slide.shapes.add_textbox(
        Inches(x + 0.3), Inches(y + 0.18), Inches(w - 0.45), Inches(h - 0.3)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    style_text(tf, title, size=title_size, bold=True, color=T.DEEP_BLUE)
    if body:
        for i, line in enumerate(body.split("\n")):
            add_paragraph(
                tf,
                line.strip(),
                size=body_size,
                color=T.CHARCOAL,
                space_before=6 if i == 0 else 3,
            )


def add_metric_box(
    slide: Slide,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    number: str,
    label: str,
    color: RGBColor = T.FOREST,
) -> None:
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.adjustments[0] = 0.12
    set_fill(box, T.WHITE)
    set_line(box, T.LIGHT_GRAY, 1)

    # Number
    num_tb = slide.shapes.add_textbox(Inches(x), Inches(y + 0.25), Inches(w), Inches(h * 0.55))
    style_text(
        num_tb.text_frame,
        number,
        size=42,
        bold=True,
        color=color,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Label
    lbl_tb = slide.shapes.add_textbox(
        Inches(x + 0.1), Inches(y + h * 0.65), Inches(w - 0.2), Inches(h * 0.3)
    )
    style_text(
        lbl_tb.text_frame,
        label,
        size=11,
        color=T.STONE,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.TOP,
    )


def add_section_label(slide: Slide, x: float, y: float, text: str) -> None:
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(6.0), Inches(0.3))
    style_text(tb.text_frame, text.upper(), size=11, bold=True, color=T.FOREST)


def add_speaker_notes(slide: Slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def new_slide(prs: Presentation) -> Slide:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    slide_background(slide, T.OFFWHITE)
    return slide


# ──────────────────────────── Slides ────────────────────────────


def build_slide_1_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, T.DEEP_BLUE)

    # Decorative accent bands
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(6.6), Inches(13.333), Inches(0.15)
    )
    set_fill(band, T.FOREST)
    set_no_line(band)
    band2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(6.8), Inches(13.333), Inches(0.06)
    )
    set_fill(band2, T.TAUPE)
    set_no_line(band2)

    # Grant code chip (top-left)
    chip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(0.6), Inches(3.3), Inches(0.45)
    )
    chip.adjustments[0] = 0.5
    set_fill(chip, T.FOREST)
    set_no_line(chip)
    chip_tb = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(3.3), Inches(0.45))
    style_text(
        chip_tb.text_frame,
        "ԲՏԱՆ-ԴՄ-2026/01  ·  ԳԱՂԱՓԱՐԻ ՓՈՒԼ",
        size=11,
        bold=True,
        color=T.WHITE,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # Main wordmark
    logo = slide.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(12), Inches(1.6))
    style_text(
        logo.text_frame,
        "HexaBin",
        size=120,
        bold=True,
        color=T.WHITE,
        align=PP_ALIGN.LEFT,
    )

    # Tagline
    tag = slide.shapes.add_textbox(Inches(0.7), Inches(3.85), Inches(12), Inches(0.7))
    style_text(
        tag.text_frame,
        "Այսօր տեսակավորիր  ·  վաղը փրկիր",
        size=28,
        color=T.TAUPE,
        align=PP_ALIGN.LEFT,
    )

    # Subtitle
    sub = slide.shapes.add_textbox(Inches(0.7), Inches(4.7), Inches(12), Inches(0.5))
    style_text(
        sub.text_frame,
        "Խելացի աղբամանների և տվյալների վաճառքի AI համակարգ",
        size=18,
        color=T.WHITE,
    )

    # Bottom metadata
    meta = slide.shapes.add_textbox(Inches(0.7), Inches(6.95), Inches(12), Inches(0.4))
    tf = meta.text_frame
    tf.word_wrap = True
    style_text(
        tf,
        "«ԷՄ ՔՅՈՒԲ» ՍՊԸ     ·     Դրամաշնորհային առաջարկ     ·     Մայիս 2026     ·     ք․ Երևան",
        size=12,
        color=T.WHITE,
    )

    add_speaker_notes(
        slide,
        "Բարև ձեզ։ Ես ներկայացնում եմ HexaBin-ը՝ արհեստական բանականությամբ կառավարվող "
        "խելացի աղբամանների համակարգ, որը թափոնների տեսակավորումը վերածում է "
        "եկամտաբեր տվյալների հոսքի։ Դիմում ենք ԲՏԱՆ-ԴՄ-2026/01 դրամաշնորհին՝ "
        "գաղափարի փուլում։",
    )


def build_slide_2_summary(prs: Presentation) -> None:
    slide = new_slide(prs)
    add_title_bar(slide, "Համառոտագիր", 2)

    # Hero sentence
    hero = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.33), Inches(1.3))
    tf = hero.text_frame
    tf.word_wrap = True
    style_text(
        tf,
        "HexaBin-ը սովորական աղբամանները վերածում է գերխելացի հանգույցների՝ 6 տեսակի "
        "թափոնի ավտոմատ տեսակավորման և սպառողական վարքագծի իրական-ժամանակյա "
        "տվյալների բազայի, որը մենք վաճառում ենք FMCG կորպորացիաներին։",
        size=18,
        color=T.CHARCOAL,
    )

    # 4 metric boxes
    y = 3.6
    h = 2.3
    w = 2.85
    gap = 0.3
    total_w = 4 * w + 3 * gap
    start_x = (13.333 - total_w) / 2

    specs = [
        ("6", "Աղբի տեսակ՝ տեղում սորտավորված", T.FOREST),
        ("10,000+", "Նկարների դատասեթ մեկ տարում", T.DEEP_BLUE),
        ("25", "Ակտիվ հանգույց 18 ամսում", T.TAUPE),
        ("3", "Եկամտի անկախ հոսք", T.ORGANIC),
    ]
    for i, (num, lbl, color) in enumerate(specs):
        add_metric_box(
            slide,
            start_x + i * (w + gap),
            y,
            w,
            h,
            number=num,
            label=lbl,
            color=color,
        )

    add_footer(slide)
    add_speaker_notes(
        slide,
        "HexaBin-ը լուծում է երկու խնդիր միաժամանակ՝ ֆիզիկական աղբատեսակավորում և "
        "թվային վերլուծական շերտ։ Մեկուկես տարում թիրախում ենք 25 ակտիվ հանգույց "
        "Երևանում, 10,000+ նկարից բաղկացած հայկական դատասեթ և երեք անկախ "
        "եկամտի հոսք։",
    )


def _draw_progress_bar(
    slide: Slide,
    x: float,
    y: float,
    w: float,
    pct: float,
    accent: RGBColor,
) -> None:
    """Horizontal progress bar: grey background + accent fill to pct percent."""
    h = 0.3
    # Background track
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    bg.adjustments[0] = 0.5
    set_fill(bg, T.LIGHT_GRAY)
    set_no_line(bg)
    # Filled portion
    fill_w = max(0.1, w * pct / 100.0)
    fill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(fill_w), Inches(h)
    )
    fill.adjustments[0] = 0.5
    set_fill(fill, accent)
    set_no_line(fill)


def build_slide_3_problem(prs: Presentation) -> None:
    slide = new_slide(prs)
    add_title_bar(slide, "Երեք խնդիր, որ մենք լուծում ենք", 3)

    add_section_label(slide, 0.5, 1.1, "ԻՐԱՎԻՃԱԿԸ ԱՅՍՕՐ՝ ԹՎԵՐՈՎ")

    card_y = 1.6
    card_h = 5.4
    card_w = 4.1
    gap = 0.15
    start_x = (13.333 - 3 * card_w - 2 * gap) / 2

    problems = [
        {
            "title": "Սորտավորման բացակայություն",
            "body": "Հայաստանում գործնականում չկա ավտոմատ աղբասորտավորում։ "
                    "Արժեքավոր հումքը կորչում է խառը թափոնների հետ։",
            "stat_num": "95%",
            "stat_pct": 95,
            "stat_caption": "Հայաստանում արտադրվող աղբից չի սորտավորվում",
            "source": "Աղբյուր՝ UN/World Bank գնահատականներ",
            "accent": T.ERROR,
        },
        {
            "title": "«Կույր» աղբահանություն",
            "body": "Մեքենաները շրջում են հաստատված գրաֆիկով՝ անկախ աղբամանի "
                    "լցվածությունից։ Վառելիք, ժամանակ ու էկոլոգիա՝ վատնված։",
            "stat_num": "30%",
            "stat_pct": 30,
            "stat_caption": "Վառելիքի և մարդկային ռեսուրսի կորուստ ավանդական ցիկլում",
            "source": "Աղբյուր՝ Smart Waste Management ուսումնասիրություններ",
            "accent": T.WARNING,
        },
        {
            "title": "FMCG-ը չունի սպառման տվյալներ",
            "body": "Ապրանքանիշները միլիոններ են ծախսում մարքեթինգի վրա, "
                    "բայց չունեն ֆիզիկական տվյալ՝ ՈՐՏԵՂ և ԵՐԲ է սպառվել ապրանքը։",
            "stat_num": "$85 ՄԼՐԴ",
            "stat_pct": 100,
            "stat_caption": "Մարքեթինգային հետազոտության համաշխարհային ծախս՝ "
                            "առանց real-world disposal data-ի",
            "source": "Աղբյուր՝ ESOMAR Global Market Research 2024",
            "accent": T.DEEP_BLUE,
        },
    ]

    for i, p in enumerate(problems):
        x = start_x + i * (card_w + gap)
        # Card shell
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(card_y), Inches(card_w), Inches(card_h),
        )
        card.adjustments[0] = 0.05
        set_fill(card, T.WHITE)
        set_line(card, T.LIGHT_GRAY, 1)

        # Top accent band
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(card_y), Inches(card_w), Inches(0.12),
        )
        set_fill(band, p["accent"])
        set_no_line(band)

        # Title
        title_tb = slide.shapes.add_textbox(
            Inches(x + 0.25), Inches(card_y + 0.25), Inches(card_w - 0.5), Inches(0.8)
        )
        ttf = title_tb.text_frame
        ttf.word_wrap = True
        style_text(ttf, p["title"], size=17, bold=True, color=T.DEEP_BLUE)

        # Body
        body_tb = slide.shapes.add_textbox(
            Inches(x + 0.25), Inches(card_y + 1.15), Inches(card_w - 0.5), Inches(1.5)
        )
        btf = body_tb.text_frame
        btf.word_wrap = True
        style_text(btf, p["body"], size=12, color=T.CHARCOAL)

        # Chart area divider
        div = slide.shapes.add_connector(
            1, Inches(x + 0.25), Inches(card_y + 2.65),
            Inches(x + card_w - 0.25), Inches(card_y + 2.65),
        )
        div.line.color.rgb = T.LIGHT_GRAY
        div.line.width = Pt(0.5)

        # Big number
        num_tb = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(card_y + 2.85), Inches(card_w - 0.3), Inches(1.1)
        )
        style_text(
            num_tb.text_frame, p["stat_num"],
            size=50, bold=True, color=p["accent"],
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Progress bar
        _draw_progress_bar(
            slide, x + 0.3, card_y + 4.0, card_w - 0.6,
            p["stat_pct"], p["accent"],
        )

        # Caption
        cap_tb = slide.shapes.add_textbox(
            Inches(x + 0.25), Inches(card_y + 4.4), Inches(card_w - 0.5), Inches(0.7)
        )
        ctf = cap_tb.text_frame
        ctf.word_wrap = True
        style_text(
            ctf, p["stat_caption"],
            size=11, color=T.CHARCOAL, align=PP_ALIGN.CENTER,
        )

        # Source
        src_tb = slide.shapes.add_textbox(
            Inches(x + 0.25), Inches(card_y + card_h - 0.5), Inches(card_w - 0.5), Inches(0.3)
        )
        style_text(
            src_tb.text_frame, p["source"],
            size=9, color=T.STONE, align=PP_ALIGN.CENTER,
        )

    add_footer(slide)
    add_speaker_notes(
        slide,
        "Երեք խնդիրը՝ Հայաստանում աղբի 95%-ը չի սորտավորվում, ավանդական "
        "աղբահանությունը կորցնում է ~30% արդյունավետություն, և FMCG-ները "
        "ծախսում են $85 միլիարդ դոլար մարքեթինգի վրա առանց real-world "
        "disposal data-ի։ HexaBin-ը լուծում է երեքը միաժամանակ։",
    )


def build_slide_4_solution(prs: Presentation) -> None:
    slide = new_slide(prs)
    add_title_bar(slide, "Մեկ սարք՝ ֆիզիկական ու թվային լուծումով", 4)

    add_section_label(slide, 0.5, 1.05, "ՄԵՐ ԼՈՒԾՈՒՄԸ")

    # Intro synthesis statement
    intro = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.45))
    style_text(
        intro.text_frame,
        "HexaBin-ը միավորում է ֆիզիկական սորտավորումն ու թվային վերլուծությունը "
        "մեկ հանգույցում՝ վերածելով աղբը եկամտի։",
        size=13, color=T.CHARCOAL, align=PP_ALIGN.CENTER,
    )

    layers = [
        {
            "title": "🧱  ՖԻԶԻԿԱԿԱՆ ՇԵՐՏ",
            "subtitle": "Տեղում ավտոմատ սորտավորում",
            "accent": T.FOREST,
            "x": 0.5,
            "features": [
                ("🎯", "6-խցիկ ավտոմատ սորտավորում",
                 "Անկախ խցիկներ՝ յուրաքանչյուր աղբի տեսակի համար"),
                ("👁️", "Խաչաձև ստուգման սենսորներ",
                 "Camera · Ultrasonic ToF · Load Cell"),
                ("⚙️", "Ճշգրիտ շարժիչներ",
                 "Սերվո + քայլային շարժիչներ՝ բաժանման համար"),
                ("🛡️", "Դիմացկուն պատյան",
                 "Չժանգոտվող պողպատ · անխոցելի վանդալիզմին"),
            ],
        },
        {
            "title": "📡  ԹՎԱՅԻՆ ՇԵՐՏ",
            "subtitle": "Ամպային վերլուծություն և տվյալներ",
            "accent": T.DEEP_BLUE,
            "x": 7.13,
            "features": [
                ("🌡️", "Միկրոկլիմայի տելեմետրիա",
                 "Գազ · ջերմություն · աղմուկ · խոնավություն"),
                ("☁️", "Edge → Cloud ML",
                 "Շարունակական վարժեցում AWS GPU-ի վրա"),
                ("🗺️", "Խելացի երթուղիացում",
                 "−30% վառելիք աղբահան մեքենաների համար"),
                ("💼", "B2B Data API",
                 "FMCG-ին վաճառվող սպառման տվյալներ"),
            ],
        },
    ]

    panel_w = 5.7
    panel_h = 4.2
    panel_y = 1.85

    for layer in layers:
        x = layer["x"]

        # Panel shell
        panel = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(panel_y),
            Inches(panel_w), Inches(panel_h),
        )
        panel.adjustments[0] = 0.04
        set_fill(panel, T.WHITE)
        set_line(panel, T.LIGHT_GRAY, 1)

        # Header bar
        hdr = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x + 0.04), Inches(panel_y),
            Inches(panel_w - 0.08), Inches(0.85),
        )
        set_fill(hdr, layer["accent"])
        set_no_line(hdr)

        hdr_tb = slide.shapes.add_textbox(
            Inches(x + 0.25), Inches(panel_y + 0.08), Inches(panel_w - 0.5), Inches(0.75),
        )
        htf = hdr_tb.text_frame
        htf.word_wrap = True
        style_text(
            htf, layer["title"], size=18, bold=True, color=T.WHITE,
            align=PP_ALIGN.LEFT,
        )
        add_paragraph(
            htf, layer["subtitle"], size=10, color=T.WHITE,
            align=PP_ALIGN.LEFT, space_before=0,
        )

        # Feature tiles
        tiles_top = panel_y + 1.0
        tile_h = 0.72
        tile_gap = 0.08
        tile_x = x + 0.2
        tile_w = panel_w - 0.4

        for i, (icon, title, desc) in enumerate(layer["features"]):
            ty = tiles_top + i * (tile_h + tile_gap)

            if i > 0:
                div = slide.shapes.add_connector(
                    1, Inches(tile_x + 0.85), Inches(ty - tile_gap / 2),
                    Inches(tile_x + tile_w), Inches(ty - tile_gap / 2),
                )
                div.line.color.rgb = T.LIGHT_GRAY
                div.line.width = Pt(0.5)

            # Icon circle on left
            icon_dia = 0.55
            icon_bg = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(tile_x + 0.05), Inches(ty + (tile_h - icon_dia) / 2),
                Inches(icon_dia), Inches(icon_dia),
            )
            set_fill(icon_bg, layer["accent"])
            set_no_line(icon_bg)
            icon_tb = slide.shapes.add_textbox(
                Inches(tile_x + 0.05), Inches(ty + (tile_h - icon_dia) / 2),
                Inches(icon_dia), Inches(icon_dia),
            )
            style_text(
                icon_tb.text_frame, icon, size=18, color=T.WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )

            # Title + description on right
            text_tb = slide.shapes.add_textbox(
                Inches(tile_x + 0.8), Inches(ty + 0.05),
                Inches(tile_w - 0.85), Inches(tile_h - 0.1),
            )
            ttf = text_tb.text_frame
            ttf.word_wrap = True
            style_text(ttf, title, size=12, bold=True, color=T.DEEP_BLUE)
            add_paragraph(ttf, desc, size=10, color=T.CHARCOAL, space_before=2)

    # Center "+" badge between panels
    plus_dia = 0.65
    plus_x = (13.333 - plus_dia) / 2
    plus_y = panel_y + panel_h / 2 - plus_dia / 2
    plus = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(plus_x), Inches(plus_y),
        Inches(plus_dia), Inches(plus_dia),
    )
    set_fill(plus, T.TAUPE)
    set_no_line(plus)
    plus_tb = slide.shapes.add_textbox(
        Inches(plus_x), Inches(plus_y), Inches(plus_dia), Inches(plus_dia),
    )
    style_text(
        plus_tb.text_frame, "+", size=30, bold=True, color=T.WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    # Outcome banner
    out_y = panel_y + panel_h + 0.2
    outcome = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(out_y),
        Inches(12.33), Inches(0.65),
    )
    outcome.adjustments[0] = 0.35
    set_fill(outcome, T.FOREST)
    set_no_line(outcome)

    eq_tb = slide.shapes.add_textbox(
        Inches(0.9), Inches(out_y), Inches(1.5), Inches(0.65),
    )
    style_text(
        eq_tb.text_frame, "= ԱՐԴՅՈՒՆՔ", size=13, bold=True, color=T.TAUPE,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )
    out_tb = slide.shapes.add_textbox(
        Inches(2.6), Inches(out_y), Inches(10.2), Inches(0.65),
    )
    style_text(
        out_tb.text_frame,
        "💧 Մաքուր հումք    ·    🚛 Օպտիմալացված լոգիստիկա    ·    📊 Տվյալների վաճառք",
        size=13, bold=True, color=T.WHITE,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )

    add_footer(slide)
    add_speaker_notes(
        slide,
        "HexaBin-ը միացնում է ֆիզիկական շերտը (6-խցիկ ավտոմատ սորտավորում՝ "
        "խաչաձև սենսորներով ու դիմացկուն պատյանով) և թվային շերտը "
        "(միկրոկլիմայի տելեմետրիա, edge-to-cloud ML, խելացի երթուղիացում "
        "և B2B Data API)։ Արդյունքը՝ մաքուր հումք, օպտիմալացված լոգիստիկա "
        "և վաճառելի տվյալներ։",
    )


def build_slide_5_architecture(prs: Presentation) -> None:
    slide = new_slide(prs)
    add_title_bar(slide, "Ինչպե՞ս է աշխատում համակարգը", 5)

    add_section_label(slide, 0.5, 1.05, "ՏՎՅԱԼՆԵՐԻ ՀՈՍՔԸ՝ ՉՈՐՍ ՓՈՒԼ")

    steps = [
        ("01", "ՀԱՎԱՔՈՒՄ", "Տեսախցիկ + Սենսոր",
         ["Dual OAK-D Stereo", "ToF · Load Cell", "Gas · Temp · IMU"]),
        ("02", "ՎԵՐԼՈՒԾՈՒՄ", "Edge AI",
         ["Jetson Orin Nano", "NN դասակարգում", "< 200 ms հետաձգում"]),
        ("03", "ՈՒՍՈՒՑՈՒՄ", "Ամպ",
         ["AWS GPU վարժեցում", "PostgreSQL պահեստ", "CI/CD թարմացում"]),
        ("04", "ՄԱՏԱԿԱՐԱՐՈՒՄ", "Արդյունք",
         ["Քաղաքային վահանակ", "Data API (B2B)", "Վերլուծ. հաշվետ."]),
    ]

    # Soft off-white zone panel that contains the 4 phase cards
    panel_x = 0.5
    panel_y = 1.5
    panel_w = 12.33
    panel_h = 3.75

    panel = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(panel_x), Inches(panel_y),
        Inches(panel_w), Inches(panel_h),
    )
    panel.adjustments[0] = 0.02
    set_fill(panel, T.OFFWHITE)
    set_line(panel, T.LIGHT_GRAY, 0.5)

    # Layout math for the 4 white cards inside the panel
    pad = 0.25
    inner_gap = 0.18
    card_w = (panel_w - 2 * pad - 3 * inner_gap) / 4  # ≈ 2.82 in
    card_h = panel_h - 2 * pad
    card_y = panel_y + pad

    for i, (num, tag, title, bullets) in enumerate(steps):
        x = panel_x + pad + i * (card_w + inner_gap)

        # White card with a thin light-gray border
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(card_y),
            Inches(card_w), Inches(card_h),
        )
        card.adjustments[0] = 0.04
        set_fill(card, T.WHITE)
        set_line(card, T.LIGHT_GRAY, 0.5)

        # Large Deep-Blue step number — top-left anchor
        num_tb = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(card_y + 0.22),
            Inches(1.2), Inches(0.7),
        )
        style_text(
            num_tb.text_frame, num, size=34, bold=True, color=T.DEEP_BLUE,
            align=PP_ALIGN.LEFT,
        )

        # Small forest-green accent bar — visual indicator between number and tag
        accent = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x + 0.3), Inches(card_y + 1.02),
            Inches(0.45), Inches(0.045),
        )
        set_fill(accent, T.FOREST)
        set_no_line(accent)

        # Stage tag in Forest Green (small, all-caps)
        tag_tb = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(card_y + 1.13),
            Inches(card_w - 0.6), Inches(0.28),
        )
        style_text(
            tag_tb.text_frame, tag, size=10, bold=True, color=T.FOREST,
            align=PP_ALIGN.LEFT,
        )

        # Stage title in Charcoal (the primary label — "Edge AI" etc.)
        title_tb = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(card_y + 1.43),
            Inches(card_w - 0.6), Inches(0.45),
        )
        style_text(
            title_tb.text_frame, title, size=16, bold=True, color=T.CHARCOAL,
            align=PP_ALIGN.LEFT,
        )

        # Hairline divider under the title
        div_y = card_y + 1.92
        div = slide.shapes.add_connector(
            1, Inches(x + 0.3), Inches(div_y),
            Inches(x + card_w - 0.3), Inches(div_y),
        )
        div.line.color.rgb = T.LIGHT_GRAY
        div.line.width = Pt(0.5)

        # Three spec lines in Stone Gray — left-aligned for editorial rhythm
        bullet_tb = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(div_y + 0.14),
            Inches(card_w - 0.6), Inches(1.3),
        )
        bf = bullet_tb.text_frame
        bf.word_wrap = True
        style_text(bf, bullets[0], size=11, color=T.STONE, align=PP_ALIGN.LEFT)
        for line in bullets[1:]:
            add_paragraph(bf, line, size=11, color=T.STONE,
                          align=PP_ALIGN.LEFT, space_before=6)

    # Pull-quote with Forest Green left accent bar
    cap_y = 5.5
    cap_accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(cap_y + 0.08),
        Inches(0.06), Inches(0.4),
    )
    set_fill(cap_accent, T.FOREST)
    set_no_line(cap_accent)

    cap_tb = slide.shapes.add_textbox(
        Inches(0.72), Inches(cap_y), Inches(12.1), Inches(0.55),
    )
    style_text(
        cap_tb.text_frame,
        "Յուրաքանչյուր նետված աղբ միաժամանակ դառնում է մաքուր հումք "
        "ԵՎ մեկ տող մեր տվյալների բազայում։",
        size=13, color=T.CHARCOAL,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )

    # Minimal KPI row — Deep Blue top rule + hairline separators
    m_y = 6.2
    metrics = [
        ("< 200 ms", "Դասակարգման հետաձգում"),
        ("95%+", "Թիրախային ճշգրտություն"),
        ("24/7", "Ավտոնոմ աշխատանք"),
    ]
    m_area_x = 0.5
    m_area_w = 12.33
    m_w = m_area_w / 3

    # Top rule — a slightly weightier Deep Blue line ties the KPI row to the deck theme
    top_line = slide.shapes.add_connector(
        1, Inches(m_area_x), Inches(m_y),
        Inches(m_area_x + m_area_w), Inches(m_y),
    )
    top_line.line.color.rgb = T.DEEP_BLUE
    top_line.line.width = Pt(1.25)

    for i, (v, k) in enumerate(metrics):
        mx = m_area_x + i * m_w

        if i > 0:
            sep = slide.shapes.add_connector(
                1, Inches(mx), Inches(m_y + 0.15),
                Inches(mx), Inches(m_y + 0.78),
            )
            sep.line.color.rgb = T.LIGHT_GRAY
            sep.line.width = Pt(0.5)

        v_tb = slide.shapes.add_textbox(
            Inches(mx + 0.35), Inches(m_y + 0.15),
            Inches(m_w - 0.6), Inches(0.45),
        )
        style_text(
            v_tb.text_frame, v, size=22, bold=True, color=T.DEEP_BLUE,
            align=PP_ALIGN.LEFT,
        )
        k_tb = slide.shapes.add_textbox(
            Inches(mx + 0.35), Inches(m_y + 0.58),
            Inches(m_w - 0.6), Inches(0.3),
        )
        style_text(
            k_tb.text_frame, k, size=10, color=T.STONE,
            align=PP_ALIGN.LEFT,
        )

    add_footer(slide)
    add_speaker_notes(
        slide,
        "Տվյալները սկսում են իրենց ճանապարհը տեսախցիկից ու սենսորներից, անցնում "
        "Jetson Orin Edge AI-ով, ուղարկվում ամպ՝ շարունակական վարժեցման համար, "
        "և հասնում են մեր հաճախորդներին երկու արտադրանքի տեսքով՝ "
        "քաղաքային վահանակ ու Data API։",
    )


def build_slide_6_innovation(prs: Presentation) -> None:
    slide = new_slide(prs)
    add_title_bar(slide, "Մեր նորարարության 4 սյունը", 6)

    add_section_label(slide, 0.5, 1.1, "ՏԱՐԱԾԱՇՐՋԱՆՈՒՄ ԱՆԱԼՈԳ ՉՈՒՆԵՑՈՂ ԿԱՐՈՂՈՒԹՅՈՒՆՆԵՐ")

    pillars = [
        (
            "01", "🤖",
            "Hardware + AI\nտեղում",
            "6 տեսակի ավտոմատ սորտավորում անմիջապես աղբամանի ներսում՝ "
            "առանց կենտրոնացված գործարանի կարիքի։",
            T.FOREST,
        ),
        (
            "02", "📈",
            "Շարունակական\nուսուցում",
            "Յուրաքանչյուր նետված աղբ լուսանկարվում ու ավելացվում է "
            "դատասեթում՝ ամեն շաբաթ բարելավելով մոդելը։",
            T.DEEP_BLUE,
        ),
        (
            "03", "💰",
            "Տվյալների\nառևտրայնացում",
            "Թափոնները վերածվում են մարքեթինգային ինֆորմացիայի՝ "
            "վաճառվում են FMCG և մրցակից բրենդներին։",
            T.TAUPE,
        ),
        (
            "04", "🌬️",
            "Ինտեգրված IoT\nտելեմետրիա",
            "Միևնույն սարքը չափում է օդի որակը, ջերմությունը, աղմուկը՝ "
            "ստեղծելով քաղաքի էկոլոգիական քարտեզը։",
            T.ORGANIC,
        ),
    ]

    card_w = 3.0
    card_h = 4.75
    gap = 0.11
    start_y = 1.5
    start_x = (13.333 - 4 * card_w - 3 * gap) / 2

    for i, (num, icon, title, body, accent) in enumerate(pillars):
        x = start_x + i * (card_w + gap)

        # Card shell
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(start_y),
            Inches(card_w), Inches(card_h),
        )
        card.adjustments[0] = 0.05
        set_fill(card, T.WHITE)
        set_line(card, T.LIGHT_GRAY, 1)

        # Top accent strip
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x + 0.05), Inches(start_y),
            Inches(card_w - 0.1), Inches(0.12),
        )
        set_fill(strip, accent)
        set_no_line(strip)

        # Big phase number
        num_tb = slide.shapes.add_textbox(
            Inches(x), Inches(start_y + 0.3), Inches(card_w), Inches(1.1),
        )
        style_text(
            num_tb.text_frame, num, size=64, bold=True, color=accent,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Icon
        icon_tb = slide.shapes.add_textbox(
            Inches(x), Inches(start_y + 1.45), Inches(card_w), Inches(0.6),
        )
        style_text(
            icon_tb.text_frame, icon, size=32, color=T.CHARCOAL,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Title (2 lines)
        title_tb = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(start_y + 2.15), Inches(card_w - 0.3), Inches(0.9),
        )
        ttf = title_tb.text_frame
        ttf.word_wrap = True
        lines = title.split("\n")
        style_text(ttf, lines[0], size=17, bold=True, color=accent, align=PP_ALIGN.CENTER)
        for line in lines[1:]:
            add_paragraph(ttf, line, size=17, bold=True, color=accent,
                          align=PP_ALIGN.CENTER, space_before=0)

        # Divider
        div = slide.shapes.add_connector(
            1, Inches(x + 0.6), Inches(start_y + 3.15),
            Inches(x + card_w - 0.6), Inches(start_y + 3.15),
        )
        div.line.color.rgb = accent
        div.line.width = Pt(1.5)

        # Body description
        body_tb = slide.shapes.add_textbox(
            Inches(x + 0.25), Inches(start_y + 3.3), Inches(card_w - 0.5), Inches(1.35),
        )
        btf = body_tb.text_frame
        btf.word_wrap = True
        style_text(btf, body, size=12, color=T.CHARCOAL, align=PP_ALIGN.CENTER)

    # Bottom callout
    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.6),
    )
    callout.adjustments[0] = 0.4
    set_fill(callout, T.DEEP_BLUE)
    set_no_line(callout)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(6.4), Inches(12.33), Inches(0.6))
    style_text(
        tb.text_frame,
        "Այս 4 կարողությունները մեկ սարքում համատեղող լուծում դեռ չկա։",
        size=14,
        bold=True,
        color=T.WHITE,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    add_footer(slide)
    add_speaker_notes(
        slide,
        "Այս 4 սյունը՝ տեղում AI սորտավորում, շարունակական ուսուցում, "
        "տվյալների առևտրայնացում ու IoT տելեմետրիա, միասին չունեն անալոգ "
        "Հայաստանում կամ հարակից շուկաներում։",
    )


def build_slide_7_market(prs: Presentation) -> None:
    slide = new_slide(prs)
    add_title_bar(slide, "Ում ենք սպասարկում", 7)

    add_section_label(slide, 0.5, 1.1, "ՄԵԿ ՇՈՒԿԱ, ԵՐԵՔ ՍԵԳՄԵՆՏ")

    segments = [
        (
            "B2G",
            "ՔԱՂԱՔԱՊԵՏԱՐԱՆ\nՀԱՄԱՅՆՔՆԵՐ",
            "«Խելացի քաղաք» ծրագրեր\nԵրթուղիների օպտիմալացում\nՕդի որակի մոնիթորինգ",
            "Երևանի քաղաքապետարան, մարզկենտրոններ",
            T.DEEP_BLUE,
        ),
        (
            "B2B / FMCG",
            "ԽՈՇՈՐ ԲՐԵՆԴՆԵՐ",
            "Սպառման ռեալ տվյալներ\nԺամ, վայր, կլիմա\nՄրցակիցների տվյալներ ևս",
            "Coca-Cola · Pepsi · Nestle · Danone",
            T.FOREST,
        ),
        (
            "Վերամշակողներ",
            "ԿԱՆԱՉ ԳՈՐԾԱՐԱՆՆԵՐ",
            "Մաքուր, արդեն սորտավորված հումք\nՀիմնավորված ծավալներ\nԵրկարաժամկետ մատակարարում",
            "Պլաստիկ, ալյումին, ապակի, թուղթ",
            T.TAUPE,
        ),
    ]
    card_y = 1.55
    card_h = 4.05
    card_w = 4.0
    gap = 0.22
    start_x = (13.333 - 3 * card_w - 2 * gap) / 2

    for i, (badge, title, body, examples, accent) in enumerate(segments):
        x = start_x + i * (card_w + gap)
        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(card_y), Inches(card_w), Inches(card_h)
        )
        card.adjustments[0] = 0.06
        set_fill(card, T.WHITE)
        set_line(card, T.LIGHT_GRAY, 1)

        # Accent header
        hdr = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(card_y), Inches(card_w), Inches(0.7)
        )
        set_fill(hdr, accent)
        set_no_line(hdr)
        hdr_tb = slide.shapes.add_textbox(
            Inches(x), Inches(card_y), Inches(card_w), Inches(0.7)
        )
        style_text(
            hdr_tb.text_frame, badge, size=20, bold=True, color=T.WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Subtitle
        sub_tb = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(card_y + 0.82), Inches(card_w - 0.4), Inches(0.95)
        )
        sub_tf = sub_tb.text_frame
        sub_tf.word_wrap = True
        style_text(sub_tf, title.split("\n")[0], size=13, bold=True, color=T.CHARCOAL,
                   align=PP_ALIGN.CENTER)
        for line in title.split("\n")[1:]:
            add_paragraph(sub_tf, line, size=13, bold=True, color=T.CHARCOAL,
                          align=PP_ALIGN.CENTER, space_before=2)

        # Value bullets
        body_tb = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(card_y + 1.95), Inches(card_w - 0.6), Inches(1.95)
        )
        bf = body_tb.text_frame
        bf.word_wrap = True
        lines = body.split("\n")
        style_text(bf, lines[0], size=11, color=T.CHARCOAL, bullet=True)
        for line in lines[1:]:
            add_paragraph(bf, line, size=11, color=T.CHARCOAL, space_before=5, bullet=True)

    # ── Confirmed-partnerships trust bar ──────────────────────────
    strip_x = 0.5
    strip_y = 5.85
    strip_w = 12.33
    strip_h = 1.1

    strip = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(strip_x), Inches(strip_y),
        Inches(strip_w), Inches(strip_h),
    )
    strip.adjustments[0] = 0.04
    set_fill(strip, T.OFFWHITE)
    set_line(strip, T.LIGHT_GRAY, 0.5)

    # Forest-Green left accent strip
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(strip_x), Inches(strip_y),
        Inches(0.08), Inches(strip_h),
    )
    set_fill(accent_bar, T.FOREST)
    set_no_line(accent_bar)

    # Section label
    label_tb = slide.shapes.add_textbox(
        Inches(strip_x + 0.28), Inches(strip_y + 0.12),
        Inches(strip_w - 0.5), Inches(0.26),
    )
    style_text(
        label_tb.text_frame, "ՀԱՍՏԱՏՎԱԾ ԳՈՐԾԸՆԿԵՐՆԵՐ",
        size=10, bold=True, color=T.FOREST,
    )

    # Two partnership rows: badge + names
    rows = [
        ("B2G", "Երևանի քաղաքապետարան   ·   Աշտարակի համայնքապետարան"),
        ("B2B", "«Մարիոթ» հյուրանոցային համալիր"),
    ]
    badge_w = 0.6
    badge_h = 0.3
    row_y_start = strip_y + 0.42
    row_gap = 0.34
    for ri, (b_text, names) in enumerate(rows):
        ry = row_y_start + ri * row_gap
        badge = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(strip_x + 0.28), Inches(ry),
            Inches(badge_w), Inches(badge_h),
        )
        badge.adjustments[0] = 0.25
        set_fill(badge, T.DEEP_BLUE)
        set_no_line(badge)
        b_tb = slide.shapes.add_textbox(
            Inches(strip_x + 0.28), Inches(ry), Inches(badge_w), Inches(badge_h),
        )
        style_text(
            b_tb.text_frame, b_text, size=10, bold=True, color=T.WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        names_tb = slide.shapes.add_textbox(
            Inches(strip_x + 0.28 + badge_w + 0.18), Inches(ry),
            Inches(strip_w - 0.7 - badge_w), Inches(badge_h),
        )
        style_text(
            names_tb.text_frame, names, size=12, color=T.CHARCOAL,
            anchor=MSO_ANCHOR.MIDDLE,
        )

    add_footer(slide)
    add_speaker_notes(
        slide,
        "Մեր շուկան երեք սեգմենտ ունի՝ B2G՝ քաղաքապետարաններ, B2B՝ FMCG "
        "կորպորացիաներ, և վերամշակող ընկերություններ, որոնք գնում են մաքուր "
        "հումքը։ Ամենաեկամտաբեր սեգմենտը FMCG-ն է։ Արդեն ունենք "
        "համագործակցության կապեր Երևանի քաղաքապետարանի և Աշտարակի "
        "համայնքապետարանի հետ, ինչպես նաև «Մարիոթ» հյուրանոցային "
        "համալիրի հետ։",
    )


def build_slide_8_business(prs: Presentation) -> None:
    slide = new_slide(prs)
    add_title_bar(slide, "Եկամտի 3 անկախ հոսք", 8)

    streams = [
        (
            "1",
            "Սորտավորված հումքի վաճառք",
            "Տեղում սորտավորված պլաստիկ, ապակի, թուղթ, ալյումին՝ "
            "վերամշակող գործարաններին։",
            "~ $150–250 / աղբաման / ամիս",
            T.PLASTIC,
            False,
        ),
        (
            "2",
            "Hardware + SaaS բաժանորդագրություն",
            "Սարքի վաճառք կամ վարձակալություն ($500–600) + "
            "ամսական պլատֆորմի վճար ($20–30/սարք)։",
            "MRR, կրկնվող եկամուտ",
            T.DEEP_BLUE,
            False,
        ),
        (
            "3",
            "Տվյալների առևտրայնացում",
            "FMCG-ներին վաճառում ենք ճշգրիտ սպառման վիճակագրություն՝ "
            "ինչ, որտեղ, երբ, ինչ կլիմայում։ Ներառյալ մրցակիցների տվյալները։",
            "Ամենաբարձր մարժա · 60%+",
            T.FOREST,
            True,
        ),
    ]

    # Three stacked horizontal cards
    card_x = 0.5
    card_w = 12.33
    card_h = 1.55
    gap = 0.17
    start_y = 1.25
    for i, (num, title, body, metric, accent, is_crown) in enumerate(streams):
        y = start_y + i * (card_h + gap)
        # Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(card_x), Inches(y), Inches(card_w), Inches(card_h)
        )
        card.adjustments[0] = 0.1
        if is_crown:
            set_fill(card, T.FOREST)
            set_no_line(card)
            title_color = T.WHITE
            body_color = T.WHITE
            metric_color = T.TAUPE
            num_color = T.WHITE
            num_bg = T.TAUPE
        else:
            set_fill(card, T.WHITE)
            set_line(card, T.LIGHT_GRAY, 1)
            title_color = T.DEEP_BLUE
            body_color = T.CHARCOAL
            metric_color = T.FOREST
            num_color = T.WHITE
            num_bg = accent

        # Number circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(card_x + 0.3), Inches(y + 0.3),
            Inches(0.95), Inches(0.95),
        )
        set_fill(circle, num_bg)
        set_no_line(circle)
        num_tb = slide.shapes.add_textbox(
            Inches(card_x + 0.3), Inches(y + 0.3), Inches(0.95), Inches(0.95)
        )
        style_text(
            num_tb.text_frame, num, size=30, bold=True, color=num_color,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Title + body
        text_tb = slide.shapes.add_textbox(
            Inches(card_x + 1.5), Inches(y + 0.2), Inches(8.0), Inches(card_h - 0.3)
        )
        tf = text_tb.text_frame
        tf.word_wrap = True
        style_text(tf, title, size=17, bold=True, color=title_color)
        add_paragraph(tf, body, size=12, color=body_color, space_before=5)

        # Metric (right)
        metric_tb = slide.shapes.add_textbox(
            Inches(card_x + 9.6), Inches(y + 0.35), Inches(2.55), Inches(card_h - 0.5)
        )
        style_text(
            metric_tb.text_frame, metric, size=14, bold=True, color=metric_color,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        if is_crown:
            badge = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(card_x + card_w - 2.2), Inches(y - 0.2),
                Inches(2.0), Inches(0.4),
            )
            badge.adjustments[0] = 0.5
            set_fill(badge, T.TAUPE)
            set_no_line(badge)
            btb = slide.shapes.add_textbox(
                Inches(card_x + card_w - 2.2), Inches(y - 0.2),
                Inches(2.0), Inches(0.4),
            )
            style_text(
                btb.text_frame, "★  ԱՄԵՆԱԵԿԱՄՏԱԲԵՐ  ★",
                size=10, bold=True, color=T.DEEP_BLUE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
            )

    # Bottom explanation
    cap = slide.shapes.add_textbox(Inches(0.5), Inches(6.35), Inches(12.33), Inches(0.6))
    style_text(
        cap.text_frame,
        "Ի տարբերություն supermarket տվյալների, մեր Data-ն չի սահմանափակվում "
        "մրցակցային արգելքներով՝ կարող ենք Coca-Cola-ին վաճառել նաև Pepsi-ի "
        "սպառման ինֆորմացիան։",
        size=12,
        color=T.STONE,
        align=PP_ALIGN.CENTER,
    )

    add_footer(slide)
    add_speaker_notes(
        slide,
        "Եկամտի երեք հոսք՝ հումքի վաճառք, hardware+SaaS բաժանորդագրություն և "
        "տվյալների վաճառք։ Երրորդը՝ Data monetization-ը, մեր ամենաբարձր մարժան "
        "ունեցող ուղղությունն է և մեր կայունության գրավականը։",
    )


def build_slide_9_competition(prs: Presentation) -> None:
    slide = new_slide(prs)
    add_title_bar(slide, "Սովորական smart bin vs HexaBin", 9)

    add_section_label(slide, 0.5, 1.1, "ՄՐՑԱԿՑԱՅԻՆ ԱՌԱՎԵԼՈՒԹՅՈՒՆ")

    # Left: typical smart bin
    left = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.6),
        Inches(6.1), Inches(5.2),
    )
    left.adjustments[0] = 0.06
    set_fill(left, T.WHITE)
    set_line(left, T.LIGHT_GRAY, 1)
    left_hdr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.6), Inches(6.1), Inches(0.7)
    )
    set_fill(left_hdr, T.STONE)
    set_no_line(left_hdr)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(6.1), Inches(0.7))
    style_text(
        tb.text_frame, "ՍՈՎՈՐԱԿԱՆ SMART BIN", size=16, bold=True, color=T.WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    # Quote
    q1 = slide.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(5.5), Inches(1.5))
    style_text(
        q1.text_frame, "«Ես լիքն եմ։»", size=30, bold=True, color=T.STONE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    # Limits list
    l1 = slide.shapes.add_textbox(Inches(0.9), Inches(4.3), Inches(5.3), Inches(2.4))
    lf = l1.text_frame
    lf.word_wrap = True
    style_text(lf, "Միայն ուլտրաձայնային սենսոր", size=13, color=T.CHARCOAL, bullet=True)
    add_paragraph(lf, "Չի տեսակավորում թափոնը", size=13, color=T.CHARCOAL, bullet=True, space_before=5)
    add_paragraph(lf, "Չի գեներացնում վերլուծական տվյալ", size=13, color=T.CHARCOAL, bullet=True, space_before=5)
    add_paragraph(lf, "Թանկ ներկրված՝ բարդ ինտեգրացիա", size=13, color=T.CHARCOAL, bullet=True, space_before=5)

    # Right: HexaBin
    right = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.73), Inches(1.6),
        Inches(6.1), Inches(5.2),
    )
    right.adjustments[0] = 0.06
    set_fill(right, T.FOREST)
    set_no_line(right)
    right_hdr = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6.73), Inches(1.6), Inches(6.1), Inches(0.7)
    )
    set_fill(right_hdr, T.DEEP_BLUE)
    set_no_line(right_hdr)
    tb2 = slide.shapes.add_textbox(Inches(6.73), Inches(1.6), Inches(6.1), Inches(0.7))
    style_text(
        tb2.text_frame, "HEXABIN", size=16, bold=True, color=T.WHITE,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )
    q2 = slide.shapes.add_textbox(Inches(6.93), Inches(2.55), Inches(5.7), Inches(1.8))
    qtf = q2.text_frame
    qtf.word_wrap = True
    style_text(
        qtf,
        "«Ինձնում 30 պլաստիկ շիշ կա։",
        size=20, bold=True, color=T.WHITE, align=PP_ALIGN.CENTER,
    )
    add_paragraph(
        qtf, "25-ը Coca-Cola են,", size=20, bold=True, color=T.WHITE,
        align=PP_ALIGN.CENTER, space_before=3,
    )
    add_paragraph(
        qtf, "ժամը 14:00-ին, 25°C շոգին։»", size=20, bold=True, color=T.TAUPE,
        align=PP_ALIGN.CENTER, space_before=3,
    )
    # Advantages
    l2 = slide.shapes.add_textbox(Inches(7.03), Inches(4.6), Inches(5.6), Inches(2.1))
    lf2 = l2.text_frame
    lf2.word_wrap = True
    style_text(lf2, "6 տեսակ ավտոմատ սորտավորում", size=13, color=T.WHITE, bullet=True)
    add_paragraph(lf2, "IoT միկրոկլիմա + AI ճանաչում", size=13, color=T.WHITE, bullet=True, space_before=5)
    add_paragraph(lf2, "Տեղական արտադրություն՝ արագ սպասարկում", size=13, color=T.WHITE, bullet=True, space_before=5)
    add_paragraph(lf2, "Հայերեն վահանակ + հայկական բրենդների դատասեթ", size=13, color=T.WHITE, bullet=True, space_before=5)

    add_footer(slide)
    add_speaker_notes(
        slide,
        "Սովորական smart bin-ը միայն ասում է «ես լիքն եմ»։ HexaBin-ը տալիս "
        "ապրանքի մակարդակով տվյալներ՝ ժամի, կլիմայի ու վայրի հետ միասին։ "
        "Սա բարձրացնում է տվյալների արժեքը տասնապատիկ։",
    )


def build_slide_10_team(prs: Presentation) -> None:
    slide = new_slide(prs)
    add_title_bar(slide, "Թիմը՝ ինժեներ + B2B վաճառք", 10)

    add_section_label(slide, 0.5, 1.1, "ՀԻՄՆԱԴԻՐՆԵՐ")

    members = [
        (
            "Արման Մովսեսյան",
            "Հիմնադիր · Գլխավոր ինժեներ · 3D դիզայներ",
            "«ԷՄ ՔՅՈՒԲ» ՍՊԸ տնօրեն · 6+ տարվա փորձ 3D և էլեկտրոնիկայի "
            "բնագավառում։ Արդեն պատրաստել է HexaBin-ի ֆիզիկական 3D մոդելն "
            "ու համակարգի ամբողջական ճարտարապետությունը։",
            "ՀՌՀ՝ Էլեկտրոնիկա և նանոէլեկտրոնիկա\n"
            "ThinkLab ինժեներական լաբորատորիայի համահիմնադիր",
            "ԱՄ",
            T.FOREST,
        ),
        (
            "Վիգեն Սարիբեկյան",
            "B2B վաճառք · Բիզնեսի զարգացում",
            "4+ տարվա փորձ B2B և կորպորատիվ վաճառքներում։ Կառուցում է "
            "պայմանագրային հարաբերություններ FMCG բրենդների հետ՝ "
            "մեր Data-արտադրանքի առևտրայնացման համար։",
            "ՀՌՀ՝ ընթացիկ ուսանող\n"
            "Կորպորատիվ հաճախորդների հետ աշխատանքի փորձ",
            "ՎՍ",
            T.DEEP_BLUE,
        ),
    ]
    card_y = 1.6
    card_h = 5.3
    card_w = 6.1
    gap = 0.13
    start_x = (13.333 - 2 * card_w - gap) / 2

    for i, (name, role, bio, education, initials, accent) in enumerate(members):
        x = start_x + i * (card_w + gap)
        # Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(card_y),
            Inches(card_w), Inches(card_h),
        )
        card.adjustments[0] = 0.06
        set_fill(card, T.WHITE)
        set_line(card, T.LIGHT_GRAY, 1)

        # Avatar circle
        avatar = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.4), Inches(card_y + 0.4),
            Inches(1.6), Inches(1.6),
        )
        set_fill(avatar, accent)
        set_no_line(avatar)
        init_tb = slide.shapes.add_textbox(
            Inches(x + 0.4), Inches(card_y + 0.4), Inches(1.6), Inches(1.6)
        )
        style_text(
            init_tb.text_frame, initials, size=50, bold=True, color=T.WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Name + role
        name_tb = slide.shapes.add_textbox(
            Inches(x + 2.2), Inches(card_y + 0.45), Inches(card_w - 2.4), Inches(1.5)
        )
        ntf = name_tb.text_frame
        ntf.word_wrap = True
        style_text(ntf, name, size=22, bold=True, color=T.DEEP_BLUE)
        add_paragraph(ntf, role, size=12, color=T.FOREST, bold=True, space_before=3)

        # Bio
        bio_tb = slide.shapes.add_textbox(
            Inches(x + 0.4), Inches(card_y + 2.3), Inches(card_w - 0.8), Inches(1.8)
        )
        btf = bio_tb.text_frame
        btf.word_wrap = True
        style_text(btf, bio, size=12, color=T.CHARCOAL)

        # Education
        edu_tb = slide.shapes.add_textbox(
            Inches(x + 0.4), Inches(card_y + 4.25), Inches(card_w - 0.8), Inches(0.95)
        )
        etf = edu_tb.text_frame
        etf.word_wrap = True
        style_text(etf, "ԿՐԹՈՒԹՅՈՒՆ", size=10, bold=True, color=T.FOREST)
        for j, line in enumerate(education.split("\n")):
            add_paragraph(etf, line, size=11, color=T.STONE, space_before=3 if j == 0 else 2)

    add_footer(slide)
    add_speaker_notes(
        slide,
        "Թիմը համատեղում է ինժեներական խորությունը և բիզնեսի զարգացումը։ "
        "Արմանն ունի 6+ տարվա փորձ 3D և էլեկտրոնիկայի ոլորտում, "
        "Վիգենը ղեկավարում է FMCG կորպորացիաների հետ պայմանագրերի կնքումը։",
    )


def build_slide_11_roadmap(prs: Presentation) -> None:
    slide = new_slide(prs)
    add_title_bar(slide, "Ճանապարհային քարտեզ · 12 ամիս", 11)

    add_section_label(slide, 0.5, 1.1, "5 ՀԱՋՈՐԴԱԿԱՆ ՓՈՒԼ · 12 ԱՄԻՍ")

    phases = [
        ("01", "Ամիսներ 1–3", "Նախապատրաստում", "Իրավական ձևակերպում",
         "ՍՊԸ + պատենտ", T.TAUPE),
        ("02", "Ամիսներ 3–5", "MVP կառուցում", "CNC · 3D · PCB",
         "Աշխատող նախատիպ", T.FOREST),
        ("03", "Ամիսներ 4–7", "Software + AI", "Cloud setup",
         "Edge-to-cloud pipeline", T.DEEP_BLUE),
        ("04", "Ամիսներ 7–10", "Պիլոտային ծրագիր", "Օպտիմալացում",
         "5 տեղադրված բին", T.ORGANIC),
        ("05", "Ամիսներ 10–12", "Մասշտաբավորում", "B2B / B2G վաճառք",
         "Առաջին պայմանագիր", T.GLASS),
    ]

    card_w = 2.25
    card_h = 4.35
    arrow_w = 0.22
    total_w = 5 * card_w + 4 * arrow_w
    start_x = (13.333 - total_w) / 2
    start_y = 1.55

    for i, (num, months, title, subtitle, deliverable, accent) in enumerate(phases):
        x = start_x + i * (card_w + arrow_w)

        # Card shell
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(start_y),
            Inches(card_w), Inches(card_h),
        )
        card.adjustments[0] = 0.06
        set_fill(card, T.WHITE)
        set_line(card, T.LIGHT_GRAY, 1)

        # Top accent strip
        strip = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x + 0.05), Inches(start_y),
            Inches(card_w - 0.1), Inches(0.1),
        )
        set_fill(strip, accent)
        set_no_line(strip)

        # Big phase number
        num_tb = slide.shapes.add_textbox(
            Inches(x), Inches(start_y + 0.2), Inches(card_w), Inches(1.0),
        )
        style_text(
            num_tb.text_frame, num, size=58, bold=True, color=accent,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Month range pill
        mo_pill = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.25), Inches(start_y + 1.3),
            Inches(card_w - 0.5), Inches(0.38),
        )
        mo_pill.adjustments[0] = 0.5
        set_fill(mo_pill, accent)
        set_no_line(mo_pill)
        mo_tb = slide.shapes.add_textbox(
            Inches(x + 0.25), Inches(start_y + 1.3), Inches(card_w - 0.5), Inches(0.38),
        )
        style_text(
            mo_tb.text_frame, months, size=11, bold=True, color=T.WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Title
        title_tb = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(start_y + 1.85), Inches(card_w - 0.2), Inches(0.4),
        )
        ttf = title_tb.text_frame
        ttf.word_wrap = True
        style_text(
            ttf, title, size=14, bold=True, color=T.DEEP_BLUE,
            align=PP_ALIGN.CENTER,
        )
        # Subtitle
        sub_tb = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(start_y + 2.3), Inches(card_w - 0.2), Inches(0.35),
        )
        sf = sub_tb.text_frame
        sf.word_wrap = True
        style_text(
            sf, subtitle, size=10, color=T.STONE,
            align=PP_ALIGN.CENTER,
        )

        # Divider
        div = slide.shapes.add_connector(
            1, Inches(x + 0.3), Inches(start_y + 2.85),
            Inches(x + card_w - 0.3), Inches(start_y + 2.85),
        )
        div.line.color.rgb = T.LIGHT_GRAY
        div.line.width = Pt(0.75)

        # Deliverable label
        lbl_tb = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(start_y + 3.0), Inches(card_w - 0.2), Inches(0.28),
        )
        style_text(
            lbl_tb.text_frame, "ԱՐԴՅՈՒՆՔ", size=9, bold=True, color=T.STONE,
            align=PP_ALIGN.CENTER,
        )
        # Deliverable
        del_tb = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(start_y + 3.35), Inches(card_w - 0.2), Inches(0.9),
        )
        dtf = del_tb.text_frame
        dtf.word_wrap = True
        style_text(
            dtf, deliverable, size=12, bold=True, color=accent,
            align=PP_ALIGN.CENTER,
        )

        # Arrow to next card
        if i < 4:
            arrow_x = x + card_w + 0.02
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(arrow_x),
                Inches(start_y + card_h / 2 - 0.15),
                Inches(arrow_w - 0.04),
                Inches(0.3),
            )
            set_fill(arrow, T.STONE)
            set_no_line(arrow)

    # Bottom 18-month target callout
    callout = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.7),
    )
    callout.adjustments[0] = 0.35
    set_fill(callout, T.FOREST)
    set_no_line(callout)
    # Left label
    lbl_tb = slide.shapes.add_textbox(
        Inches(0.9), Inches(6.2), Inches(3.0), Inches(0.7),
    )
    style_text(
        lbl_tb.text_frame, "18-ԱՄՍՅԱ ՆՊԱՏԱԿ", size=13, bold=True, color=T.TAUPE,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )
    # Right targets
    tgt_tb = slide.shapes.add_textbox(
        Inches(3.8), Inches(6.2), Inches(9.0), Inches(0.7),
    )
    style_text(
        tgt_tb.text_frame,
        "25+ ակտիվ հանգույց    ·    B2B-FMCG պայմանագիր    ·    Կայուն MRR",
        size=13, bold=True, color=T.WHITE,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
    )

    add_footer(slide)
    add_speaker_notes(
        slide,
        "Ծրագիրը բաժանված է 5 հաջորդական փուլի 12 ամսում։ Ամիսներ 1-3՝ "
        "ՍՊԸ գրանցում ու պատենտ, 3-5՝ MVP պատրաստ, 4-7՝ software + AI + "
        "cloud pipeline, 7-10՝ պիլոտ 5 բինով, 10-12՝ մասշտաբավորում և "
        "B2B/B2G առաջին պայմանագիրը։ 18 ամսում՝ 25+ ակտիվ հանգույց և "
        "կայուն MRR։",
    )


def build_slide_12_budget(prs: Presentation) -> None:
    slide = new_slide(prs)
    add_title_bar(slide, "Բյուջեն և միջոցների օգտագործումը", 12)

    add_section_label(slide, 0.5, 1.1, "ԴՐԱՄԱՇՆՈՐՀԻ ԲԱՇԽՈՒՄ")

    # Big ask on the left
    ask_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.6), Inches(5.2), Inches(5.3),
    )
    ask_card.adjustments[0] = 0.05
    set_fill(ask_card, T.DEEP_BLUE)
    set_no_line(ask_card)

    lbl_tb = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(4.8), Inches(0.4))
    style_text(
        lbl_tb.text_frame, "ԴՐԱՄԱՇՆՈՐՀԻ ԽՆԴՐԱՆՔ",
        size=12, bold=True, color=T.TAUPE,
    )
    num_tb = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(4.8), Inches(1.4))
    style_text(
        num_tb.text_frame, "5,000,000", size=62, bold=True, color=T.WHITE,
    )
    cur_tb = slide.shapes.add_textbox(Inches(0.7), Inches(3.85), Inches(4.8), Inches(0.5))
    style_text(
        cur_tb.text_frame, "ՀՀ դրամ", size=22, bold=True, color=T.WHITE,
    )
    # Divider
    div = slide.shapes.add_connector(
        1, Inches(0.7), Inches(4.7), Inches(5.5), Inches(4.7)
    )
    div.line.color.rgb = T.TAUPE
    div.line.width = Pt(1)

    # Meta info below divider
    tl = slide.shapes.add_textbox(Inches(0.7), Inches(4.9), Inches(4.8), Inches(2.0))
    tf = tl.text_frame
    tf.word_wrap = True
    style_text(tf, "ԾՐԱԳՐԻ ՄԱՆՐԱՄԱՍՆԵՐ", size=10, bold=True, color=T.TAUPE)
    add_paragraph(tf, "Ժամկետ՝ 12 ամիս", size=13, color=T.WHITE, space_before=8)
    add_paragraph(tf, "Մեկնարկ՝ 20.05.2026", size=13, color=T.WHITE, space_before=4)
    add_paragraph(tf, "Կազմակերպություն՝", size=13, color=T.WHITE, space_before=4)
    add_paragraph(tf, "«ԷՄ ՔՅՈՒԲ» ՍՊԸ", size=13, bold=True, color=T.WHITE, space_before=2)

    # Right: allocation cards (40/40/20 of 5,000,000 AMD)
    alloc = [
        ("40%", "Hardware", "Սենսորներ, PCB, 3D տպագրություն (resin), մետաղական պատյան",
         "2,000,000 դրամ", T.FOREST),
        ("40%", "Cloud + Software", "AWS GPU վարժեցում, ամպային ենթակառուցվածք, "
         "ծրագրային ապահովման մշակում", "2,000,000 դրամ", T.DEEP_BLUE),
        ("20%", "Marketing · Legal · Pilot", "ՍՊԸ գրանցում, պատենտավորում, "
         "պիլոտային թեստավորում, մարքեթինգ", "1,000,000 դրամ", T.TAUPE),
    ]
    card_x = 5.85
    card_w = 7.0
    card_h = 1.62
    gap = 0.17
    start_y = 1.6
    for i, (pct, title, body, amt, color) in enumerate(alloc):
        y = start_y + i * (card_h + gap)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(card_x), Inches(y),
            Inches(card_w), Inches(card_h),
        )
        card.adjustments[0] = 0.08
        set_fill(card, T.WHITE)
        set_line(card, T.LIGHT_GRAY, 1)

        # Percentage block
        pct_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(card_x), Inches(y),
            Inches(1.6), Inches(card_h),
        )
        set_fill(pct_bg, color)
        set_no_line(pct_bg)
        pct_tb = slide.shapes.add_textbox(
            Inches(card_x), Inches(y), Inches(1.6), Inches(card_h)
        )
        style_text(
            pct_tb.text_frame, pct, size=38, bold=True, color=T.WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )

        # Title + body
        text_tb = slide.shapes.add_textbox(
            Inches(card_x + 1.8), Inches(y + 0.2), Inches(card_w - 2.0), Inches(card_h - 0.4)
        )
        tf = text_tb.text_frame
        tf.word_wrap = True
        style_text(tf, title, size=16, bold=True, color=T.DEEP_BLUE)
        add_paragraph(tf, body, size=11, color=T.CHARCOAL, space_before=4)
        add_paragraph(tf, amt, size=11, bold=True, color=T.FOREST, space_before=4)

    add_footer(slide)
    add_speaker_notes(
        slide,
        "Մենք խնդրում ենք 5,000,000 դրամ դրամաշնորհ 12-ամսյա ծրագրի համար։ "
        "Բաշխումը՝ 40% hardware (2 մլն), 40% cloud + software (2 մլն), "
        "20% գործառնական ծախսեր (1 մլն)։",
    )


def build_slide_13_risks(prs: Presentation) -> None:
    slide = new_slide(prs)
    add_title_bar(slide, "Ռիսկերը և դրանց հաղթահարման ռազմավարություն", 13)

    risks = [
        (
            "1  ·  Տեխնոլոգիական",
            "AI ճանաչման սխալներ · մեխանիզմի խցանում",
            "Cross-validation՝ ToF + Load Cell + Camera։ "
            "10,000+ նկարից դատասեթ հայկական բրենդներով։",
            T.INFO,
        ),
        (
            "2  ·  Մատակարարման շղթա",
            "Jetson Orin-ի և հատուկ սենսորների ներկրման դեֆիցիտ",
            "Կրիտիկական էլեկտրոնիկան պատվիրվում է 1-ին ամսում։ "
            "Plan B՝ փոխարինող MCU-ներ։",
            T.WARNING,
        ),
        (
            "3  ·  Առևտրային (FMCG)",
            "Կորպորացիաների թերահավատություն Data-ի նկատմամբ",
            "1 ամսվա Free Sample Report՝ գործնականում ցույց տալու ARR "
            "և սպառման կանխատեսման ճշգրտությունը։",
            T.FOREST,
        ),
        (
            "4  ·  Ֆիզիկական անվտանգություն",
            "Վանդալիզմ · եղանակային քայքայում",
            "Չժանգոտվող պողպատ + ալյումին պատյան։ "
            "MPU-6050 աքսելերոմետր՝ ակնթարթային ահազանգ հարվածի դեպքում։",
            T.ERROR,
        ),
    ]
    card_w = 6.1
    card_h = 2.5
    gap_x = 0.13
    gap_y = 0.18
    start_x = (13.333 - 2 * card_w - gap_x) / 2
    start_y = 1.3
    for i, (title, descr, mitigation, color) in enumerate(risks):
        x = start_x + (i % 2) * (card_w + gap_x)
        y = start_y + (i // 2) * (card_h + gap_y)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
            Inches(card_w), Inches(card_h),
        )
        card.adjustments[0] = 0.06
        set_fill(card, T.WHITE)
        set_line(card, T.LIGHT_GRAY, 1)

        # Colored band (top)
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(card_w), Inches(0.5)
        )
        set_fill(band, color)
        set_no_line(band)
        btb = slide.shapes.add_textbox(
            Inches(x + 0.2), Inches(y), Inches(card_w - 0.4), Inches(0.5)
        )
        style_text(
            btb.text_frame, title, size=14, bold=True, color=T.WHITE,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # Risk description
        risk_tb = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(y + 0.65), Inches(card_w - 0.6), Inches(0.6)
        )
        style_text(risk_tb.text_frame, "ՌԻՍԿ", size=9, bold=True, color=T.ERROR)
        add_paragraph(risk_tb.text_frame, descr, size=12, color=T.CHARCOAL, space_before=2)

        # Mitigation
        mit_tb = slide.shapes.add_textbox(
            Inches(x + 0.3), Inches(y + 1.4), Inches(card_w - 0.6), Inches(card_h - 1.5)
        )
        style_text(mit_tb.text_frame, "ՀԱՂԹԱՀԱՐՈՒՄ", size=9, bold=True, color=T.FOREST)
        mf = mit_tb.text_frame
        add_paragraph(mf, mitigation, size=12, color=T.CHARCOAL, space_before=3)

    add_footer(slide)
    add_speaker_notes(
        slide,
        "Չորս հիմնական ռիսկ՝ տեխնոլոգիական, մատակարարման, առևտրային ու "
        "ֆիզիկական։ Յուրաքանչյուրի համար ունենք կոնկրետ հաղթահարման "
        "ռազմավարություն, որը նկարագրված է մանրամասն փաստաթղթում։",
    )


def build_slide_14_closing(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, T.DEEP_BLUE)

    # Left band
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), Inches(7.5)
    )
    set_fill(band, T.FOREST)
    set_no_line(band)

    # Title
    t = slide.shapes.add_textbox(Inches(0.7), Inches(0.6), Inches(12), Inches(1.0))
    tf = t.text_frame
    style_text(tf, "Մեր խնդրանքը", size=42, bold=True, color=T.WHITE)

    sub = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12), Inches(0.5))
    style_text(
        sub.text_frame,
        "5 մլն դրամ դրամաշնորհ, որ փոխի հայկական waste-tech-ի կանոնները",
        size=16,
        color=T.TAUPE,
    )

    # Three impact pillars (left)
    pillars = [
        ("📊", "ՏՆՏԵՍԱԿԱՆ",
         "-30% աղբահանության լոգիստիկ ծախս\nԿայուն MRR Հայաստանից"),
        ("🌱", "ԲՆԱՊԱՀՊԱՆԱԿԱՆ",
         "Մաքուր հումքի ավելի բարձր recycle rate\n"
         "Քաղաքի օդի որակի ռեալ քարտեզ"),
        ("♻️", "ԿԱՅՈՒՆ",
         "Data-վաճառքից ֆինանսական անկախություն\n"
         "Արտահանման ներուժ հարևան շուկաներ"),
    ]
    y = 2.4
    for i, (icon, title, body) in enumerate(pillars):
        px = 0.7 + i * 4.0
        # Icon
        icon_tb = slide.shapes.add_textbox(Inches(px), Inches(y), Inches(3.7), Inches(0.6))
        style_text(icon_tb.text_frame, icon, size=28, color=T.TAUPE)
        # Title
        title_tb = slide.shapes.add_textbox(Inches(px), Inches(y + 0.7), Inches(3.7), Inches(0.4))
        style_text(title_tb.text_frame, title, size=13, bold=True, color=T.TAUPE)
        # Body
        body_tb = slide.shapes.add_textbox(Inches(px), Inches(y + 1.15), Inches(3.7), Inches(1.5))
        bf = body_tb.text_frame
        bf.word_wrap = True
        lines = body.split("\n")
        style_text(bf, lines[0], size=12, color=T.WHITE)
        for line in lines[1:]:
            add_paragraph(bf, line, size=12, color=T.WHITE, space_before=3)

    # What the money buys
    what = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7), Inches(4.95), Inches(12.0), Inches(1.35),
    )
    what.adjustments[0] = 0.08
    set_fill(what, T.FOREST)
    set_no_line(what)
    wt = slide.shapes.add_textbox(Inches(1.0), Inches(5.05), Inches(11.5), Inches(1.2))
    wtf = wt.text_frame
    wtf.word_wrap = True
    style_text(wtf, "Ի՞ՆՉ Ե ԳՆՈՒՄ ԱՅՍ ԴՐԱՄԱՇՆՈՐՀԸ", size=12, bold=True, color=T.TAUPE)
    add_paragraph(
        wtf,
        "Առաջին աշխատող HexaBin նախատիպ · 6-խցիկ մեխանիկա · AI cloud "
        "ենթակառուցվածք · պիլոտային ծրագիր 1-2 գործընկերների մոտ · "
        "հայկական բրենդների սկզբնական դատասեթ։",
        size=14, color=T.WHITE, space_before=5,
    )

    # Contact strip
    contact = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), Inches(13.333), Inches(1.0)
    )
    set_fill(contact, T.CHARCOAL)
    set_no_line(contact)
    c_tb = slide.shapes.add_textbox(Inches(0.7), Inches(6.55), Inches(12.0), Inches(0.9))
    ctf = c_tb.text_frame
    ctf.word_wrap = True
    style_text(ctf, "«ԷՄ ՔՅՈՒԲ» ՍՊԸ  ·  ք. Երևան, Վարդանանց 104, 0070",
               size=12, bold=True, color=T.WHITE)
    add_paragraph(
        ctf,
        "Արման Մովսեսյան  ·  +374 95 851 561  ·  armovs.yan@gmail.com",
        size=12, color=T.TAUPE, space_before=3,
    )

    add_speaker_notes(
        slide,
        "Ամփոփում եմ՝ HexaBin-ը ստեղծում է տնտեսական, բնապահպանական և "
        "կայունության ազդեցություն։ Ձեր 5 միլիոն դրամ դրամաշնորհը "
        "մեզ թույլ կտա կառուցել առաջին աշխատող նախատիպն ու մեկնարկել "
        "պիլոտը։ Շնորհակալություն։",
    )


def build_slide_15_thanks(prs: Presentation) -> None:
    """Closing 'thank you' billboard — split layout: hero left, contact card right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, T.DEEP_BLUE)

    # Forest-green left accent band
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), Inches(7.5)
    )
    set_fill(band, T.FOREST)
    set_no_line(band)

    # ── LEFT HALF: hero ─────────────────────────────────────────
    # HEXABIN wordmark (small, top-left of hero zone)
    wm = slide.shapes.add_textbox(Inches(0.85), Inches(1.1), Inches(7.0), Inches(0.4))
    style_text(wm.text_frame, "HEXABIN", size=14, bold=True, color=T.TAUPE)

    # Tiny Taupe rule under wordmark
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(1.5), Inches(0.7), Inches(0.04),
    )
    set_fill(rule, T.TAUPE)
    set_no_line(rule)

    # Big headline
    title_tb = slide.shapes.add_textbox(
        Inches(0.85), Inches(2.2), Inches(7.4), Inches(2.1),
    )
    style_text(
        title_tb.text_frame, "Շնորհակալություն",
        size=64, bold=True, color=T.WHITE, align=PP_ALIGN.LEFT,
    )

    # Tagline
    tag_tb = slide.shapes.add_textbox(
        Inches(0.85), Inches(4.45), Inches(7.4), Inches(1.4),
    )
    tf = tag_tb.text_frame
    tf.word_wrap = True
    style_text(
        tf, "Միասին կառուցենք խելացի,",
        size=20, color=T.TAUPE, align=PP_ALIGN.LEFT,
    )
    add_paragraph(tf, "կանաչ քաղաքը։", size=20, color=T.TAUPE, align=PP_ALIGN.LEFT, space_before=2)

    # ── RIGHT HALF: contact card ────────────────────────────────
    card_x = 8.55
    card_y = 1.1
    card_w = 4.3
    card_h = 5.0

    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(card_x), Inches(card_y),
        Inches(card_w), Inches(card_h),
    )
    card.adjustments[0] = 0.04
    set_fill(card, T.WHITE)
    set_no_line(card)

    # Forest top accent strip on card
    top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(card_x), Inches(card_y),
        Inches(card_w), Inches(0.12),
    )
    set_fill(top, T.FOREST)
    set_no_line(top)

    # "ԿԱՊ" label
    label_tb = slide.shapes.add_textbox(
        Inches(card_x + 0.4), Inches(card_y + 0.32), Inches(card_w - 0.8), Inches(0.3),
    )
    style_text(
        label_tb.text_frame, "ԿԱՊ", size=11, bold=True, color=T.FOREST,
    )

    # Founder name (large)
    name_tb = slide.shapes.add_textbox(
        Inches(card_x + 0.4), Inches(card_y + 0.72), Inches(card_w - 0.8), Inches(0.55),
    )
    style_text(
        name_tb.text_frame, "Արման Մովսեսյան",
        size=22, bold=True, color=T.DEEP_BLUE,
    )

    # Role
    role_tb = slide.shapes.add_textbox(
        Inches(card_x + 0.4), Inches(card_y + 1.27), Inches(card_w - 0.8), Inches(0.3),
    )
    style_text(
        role_tb.text_frame, "Հիմնադիր · CEO",
        size=12, color=T.STONE,
    )

    # Hairline divider
    div = slide.shapes.add_connector(
        1, Inches(card_x + 0.4), Inches(card_y + 1.75),
        Inches(card_x + card_w - 0.4), Inches(card_y + 1.75),
    )
    div.line.color.rgb = T.LIGHT_GRAY
    div.line.width = Pt(0.5)

    # Contact lines: icon-style label + value (no emojis)
    contact_rows = [
        ("ՀԵՌԱԽՈՍ", "+374 95 851 561"),
        ("ԷԼ. ՀԱՍՑԵ", "armovs.yan@gmail.com"),
        ("ԳՐԱՍԵՆՅԱԿ", "ք. Երևան, Վարդանանց 104, 0070"),
        ("ԸՆԿԵՐՈՒԹՅՈՒՆ", "«ԷՄ ՔՅՈՒԲ» ՍՊԸ"),
    ]
    row_y = card_y + 1.95
    for lbl, val in contact_rows:
        l_tb = slide.shapes.add_textbox(
            Inches(card_x + 0.4), Inches(row_y), Inches(card_w - 0.8), Inches(0.22),
        )
        style_text(l_tb.text_frame, lbl, size=9, bold=True, color=T.STONE)
        v_tb = slide.shapes.add_textbox(
            Inches(card_x + 0.4), Inches(row_y + 0.24),
            Inches(card_w - 0.8), Inches(0.32),
        )
        style_text(v_tb.text_frame, val, size=13, bold=True, color=T.CHARCOAL)
        row_y += 0.7

    # Bottom thin caption strip — grant code reminder
    foot = slide.shapes.add_textbox(Inches(0.85), Inches(6.85), Inches(12.0), Inches(0.4))
    style_text(
        foot.text_frame,
        "ԲՏԱՆ-ԴՄ-2026/01   ·   HexaBin   ·   Smart Waste AI",
        size=10, color=T.TAUPE,
    )

    add_speaker_notes(
        slide,
        "Շնորհակալություն ձեր ուշադրության համար։ Միասին կարող ենք "
        "Հայաստանը դարձնել տարածաշրջանի առաջատարը waste-tech-ի ոլորտում։ "
        "Իմ կոնտակտները ձեր առջևում են։",
    )


def build_slide_16_qa(prs: Presentation) -> None:
    """Q&A 'billboard' — stays on screen during questions; surfaces conversation topics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_background(slide, T.DEEP_BLUE)

    # Forest-green left accent band
    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), Inches(7.5)
    )
    set_fill(band, T.FOREST)
    set_no_line(band)

    # Top-left wordmark
    wm = slide.shapes.add_textbox(Inches(0.85), Inches(0.7), Inches(6.0), Inches(0.4))
    style_text(wm.text_frame, "HEXABIN", size=14, bold=True, color=T.TAUPE)

    # Eyebrow label
    eb_tb = slide.shapes.add_textbox(Inches(0.85), Inches(1.45), Inches(12.0), Inches(0.4))
    style_text(eb_tb.text_frame, "Q  &  A", size=18, bold=True, color=T.TAUPE)

    # Main headline (left-aligned, strong)
    title_tb = slide.shapes.add_textbox(
        Inches(0.85), Inches(1.95), Inches(12.0), Inches(1.2),
    )
    style_text(
        title_tb.text_frame, "Հարցեր և պատասխաններ",
        size=46, bold=True, color=T.WHITE, align=PP_ALIGN.LEFT,
    )

    # Forest divider
    div = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.85), Inches(3.05), Inches(0.9), Inches(0.05),
    )
    set_fill(div, T.FOREST)
    set_no_line(div)

    # Subline
    sub_tb = slide.shapes.add_textbox(
        Inches(0.85), Inches(3.2), Inches(12.0), Inches(0.4),
    )
    style_text(
        sub_tb.text_frame,
        "Բաց ենք քննարկելու ցանկացած թեմա — ընտրեք, որտեղից սկսենք։",
        size=15, color=T.TAUPE,
    )

    # ── 4 topic cards (conversation starters) ───────────────────
    topics = [
        ("01", "ՏԵԽՆՈԼՈԳԻԱ",
         "Hardware, AI մոդելներ, edge-cloud ճարտարապետություն"),
        ("02", "ԲԻԶՆԵՍ",
         "Եկամտի հոսքեր, unit economics, գնագոյացում"),
        ("03", "ՇՈՒԿԱ",
         "Գործընկերներ, պիլոտներ, արտահանման պլան"),
        ("04", "ԹԻՄ",
         "Կազմ, փորձառություն, դրամաշնորհի օգտագործում"),
    ]
    card_y = 4.0
    card_h = 1.85
    card_gap = 0.18
    available_w = 13.333 - 2 * 0.85
    card_w = (available_w - 3 * card_gap) / 4

    for i, (num, label, body) in enumerate(topics):
        x = 0.85 + i * (card_w + card_gap)

        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(card_y),
            Inches(card_w), Inches(card_h),
        )
        card.adjustments[0] = 0.06
        set_fill(card, T.WHITE)
        set_no_line(card)

        # Forest left accent
        a = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(x), Inches(card_y + 0.18),
            Inches(0.06), Inches(card_h - 0.36),
        )
        set_fill(a, T.FOREST)
        set_no_line(a)

        # Number
        n_tb = slide.shapes.add_textbox(
            Inches(x + 0.22), Inches(card_y + 0.18),
            Inches(card_w - 0.4), Inches(0.42),
        )
        style_text(n_tb.text_frame, num, size=20, bold=True, color=T.DEEP_BLUE)

        # Label
        l_tb = slide.shapes.add_textbox(
            Inches(x + 0.22), Inches(card_y + 0.66),
            Inches(card_w - 0.4), Inches(0.32),
        )
        style_text(l_tb.text_frame, label, size=11, bold=True, color=T.FOREST)

        # Body
        b_tb = slide.shapes.add_textbox(
            Inches(x + 0.22), Inches(card_y + 1.0),
            Inches(card_w - 0.4), Inches(0.85),
        )
        bf = b_tb.text_frame
        bf.word_wrap = True
        style_text(bf, body, size=10, color=T.CHARCOAL)

    # ── Bottom contact strip ────────────────────────────────────
    contact = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(6.5), Inches(13.333), Inches(1.0)
    )
    set_fill(contact, T.CHARCOAL)
    set_no_line(contact)
    c_tb = slide.shapes.add_textbox(Inches(0.85), Inches(6.55), Inches(12.0), Inches(0.9))
    ctf = c_tb.text_frame
    ctf.word_wrap = True
    style_text(ctf, "Արման Մովսեսյան  ·  Հիմնադիր · CEO",
               size=12, bold=True, color=T.WHITE)
    add_paragraph(
        ctf,
        "+374 95 851 561   ·   armovs.yan@gmail.com   ·   ք. Երևան, Վարդանանց 104",
        size=12, color=T.TAUPE, space_before=3,
    )

    add_speaker_notes(
        slide,
        "Հրավիրում եմ լսարանին տալ իրենց հարցերը։ Չորս թեմա ենք առանձնացրել "
        "որպես մեկնակետ՝ տեխնոլոգիա, բիզնես, շուկա, թիմ, բայց պատրաստ ենք "
        "ցանկացած այլ հարցի։",
    )


# ──────────────────────────── Main ────────────────────────────


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    builders = [
        build_slide_1_cover,
        build_slide_2_summary,
        build_slide_3_problem,
        build_slide_4_solution,
        build_slide_5_architecture,
        build_slide_6_innovation,
        build_slide_7_market,
        build_slide_8_business,
        build_slide_9_competition,
        build_slide_10_team,
        build_slide_11_roadmap,
        build_slide_12_budget,
        build_slide_13_risks,
        build_slide_14_closing,
        build_slide_15_thanks,
        build_slide_16_qa,
    ]
    for build in builders:
        build(prs)

    prs.save(str(OUTPUT_PATH))
    print(f"Wrote {OUTPUT_PATH} ({OUTPUT_PATH.stat().st_size} bytes, {len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
